# ppt_builder.py
from pathlib import Path
from pptx import Presentation
from pptx.enum.text import PP_ALIGN

def create_ppt_from_json_safe(presentation_json, output_path, templates_dir="./templates", num_slides=None):
    TEMPLATE_MAP = {
        "education": Path(templates_dir) / "education_template.pptx",
        "technology": Path(templates_dir) / "tech_template.pptx",
        "general": Path(templates_dir) / "general_template.pptx",
        "sports": Path(templates_dir) / "sports_template.pptx",
        "health": Path(templates_dir) / "health_template.pptx"
    }

    topic = str(presentation_json.get("topic_category", "general")).strip().lower()
    template_path = TEMPLATE_MAP.get(topic, TEMPLATE_MAP["general"])
    prs = Presentation(str(template_path))

    slides_data = presentation_json.get("slides", [])
    total_slides = num_slides if num_slides else len(slides_data)

    for i, slide_data in enumerate(slides_data[:total_slides]):
        title_text = slide_data.get("title", f"Slide {i+1}")
        content_points = slide_data.get("content", [])

        if i == 0 and len(prs.slides) == 1:
            slide = prs.slides[0]
        else:
            layout_index = 0 if len(content_points) == 0 else 1
            slide = prs.slides.add_slide(prs.slide_layouts[layout_index])

        if slide.shapes.title:
            slide.shapes.title.text = title_text

        if content_points:
            for shape in slide.placeholders:
                if shape.is_placeholder and hasattr(shape, "text_frame") and shape.placeholder_format.idx != 0:
                    text_frame = shape.text_frame
                    text_frame.clear()
                    for bullet in content_points:
                        if bullet.strip():
                            p = text_frame.add_paragraph()
                            p.text = bullet
                            p.rtl = True
                            p.alignment = PP_ALIGN.RIGHT

    prs.save(output_path)
    print(f"Presentation saved at: {output_path}")
    return output_path
